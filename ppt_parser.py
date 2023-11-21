from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    """
    Extracts and prints the text from each slide in a PowerPoint file.

    :param pptx_file: The path to the PowerPoint file to be processed.
    """
    prs = Presentation(pptx_file)  # Load the PowerPoint file
    for slide_number, slide in enumerate(prs.slides):
        print("Slide {}:").format(slide_number + 1)  # Print the slide number
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)  # Print the text content of each shape

def extract_images_from_pptx(pptx_file):
    """
    Extracts images from each slide in a PowerPoint file and saves them as separate files.

    :param pptx_file: The path to the PowerPoint file to be processed.
    """
    prs = Presentation(pptx_file)  # Load the PowerPoint file
    for slide_number, slide in enumerate(prs.slides):
        print("Slide {}:").format(slide_number + 1)  # Print the slide number
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # Check if the shape is an image
                image = shape.image
                # Construct a filename for each image
                image_filename = "slide-{}-shape-{}.{}".format(slide_number + 1, shape.shape_id, image.ext)
                # Save the image to the 'assets' folder
                with open("assets/" + image_filename, "wb") as f:
                    f.write(image.blob)
                print("Saved image {}").format(image_filename)  # Print confirmation of saved image

# The path to the PowerPoint file
pptx_file = 'Luna mHealth Mockup.pptx'

# Call the function to extract text
extract_text_from_pptx(pptx_file)

# Call the function to extract images
extract_images_from_pptx(pptx_file)
