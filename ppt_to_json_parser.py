import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os
import shutil
import uuid
import re

def parse_pptx_to_json(pptx_file):
    """
    Parses a PowerPoint presentation and converts it into a JSON object.

    Args:
        pptx_file (str): The file path of the PowerPoint presentation.

    Returns:
        dict: A JSON object containing the parsed data from the presentation.
    """
    try:
        # Open the PowerPoint presentation
        prs = pptx.Presentation(pptx_file)

        # Extract the presentation name from the file path
        presentation_name = os.path.basename(pptx_file).split('.')[0]

        # Create a JSON object to store the presentation data
        json_data = {'Module': presentation_name, 'dimensions': {}, 'slides': []} 

        # Retrieve the presentation dimensions in centimeters
        presentation_width = round(prs.slide_width / 360000, 4)
        presentation_height = round(prs.slide_height / 360000, 4)

        # Add the presentation dimensions to the JSON object
        json_data['dimensions']['width'] = presentation_width
        json_data['dimensions']['height'] = presentation_height

        # Create a directory to store the images
        image_dir = 'JSON_images'

        # Delete the existing image directory if it exists
        if os.path.exists(image_dir):
            shutil.rmtree(image_dir)

        os.makedirs(image_dir, exist_ok=True)

        # Iterate over the slides in the presentation
        for slide_num, slide in enumerate(prs.slides, start=1):

            # Create a JSON object to store the slide data
            slide_data = {}

            # Add the slide number to the JSON object
            slide_data['slide_number'] = slide_num

            # Add the slide title to the JSON object
            if slide.shapes.title is not None:
                slide_data['title'] = slide.shapes.title.text

            # Add the slide notes to the JSON object
            # TODO: URLs in notes are not being extracted
            if slide.has_notes_slide:
                slide_data['notes'] = slide.notes_slide.notes_text_frame.text

                # Extract clickable text from notes
                clickable_text_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
                clickable_text_matches = re.findall(clickable_text_pattern, slide_data['notes'])
                if clickable_text_matches:
                    slide_data['clickable_text'] = []
                    for match in clickable_text_matches:
                        slide_data['clickable_text'].append({'text': match[0], 'url': match[1]})

            # Add the slide elements to the JSON object
            slide_data['elements'] = []
            for shape in slide.shapes:
                element_data = {}

                # Add the element type to the JSON object
                element_data['type'] = shape.shape_type.name if isinstance(shape.shape_type, MSO_SHAPE_TYPE) else shape.shape_type

                # Add the element text to the JSON object
                if hasattr(shape, 'text'):
                    element_data['text'] = shape.text

                    # Add the text properties to the JSON object
                    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                        element_data['text_properties'] = {
                            'font_name': shape.text_frame.paragraphs[0].runs[0].font.name,
                            'font_size': shape.text_frame.paragraphs[0].runs[0].font.size,
                            'font_bold': shape.text_frame.paragraphs[0].runs[0].font.bold,
                            'font_italic': shape.text_frame.paragraphs[0].runs[0].font.italic,
                            'font_color': shape.text_frame.paragraphs[0].runs[0].font.color.rgb if shape.text_frame.paragraphs[0].runs[0].font.color.type == 'RGB' else '',
                        }

                # Add the element image path to the JSON object
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Generate a unique filename for the image
                    image_filename = f"{shape.name}.png"
                    image_subdir = os.path.join(image_dir, f'slide_{slide_num}')
                    os.makedirs(image_subdir, exist_ok=True)
                    image_path = os.path.join(image_subdir, image_filename)
                    with open(image_path, 'wb') as f:
                        f.write(shape.image.blob)

                    # Add the image path to the JSON object 
                    element_data['image_path'] = image_path

                    # Add common properties or attributes associated with PowerPoint image objects
                    element_data['properties'] = {
                        'name': shape.name,
                        'description': shape.description if hasattr(shape, 'description') else '',
                        'extension': shape.image.ext,
                        'content_type': shape.image.content_type,
                        'file_size': shape.image.size,
                        'dpi': shape.image.dpi,
                    }

                # Add the element position in centimeters to the JSON object
                element_data['position'] = {
                    'left': round(shape.left / 360000, 4),
                    'top': round(shape.top / 360000, 4),
                    'width': round(shape.width / 360000, 4),
                    'height': round(shape.height / 360000, 4)
                }

                # Add the element data to the slide data
                slide_data['elements'].append(element_data)

            # Add the slide data to the presentation JSON object
            json_data['slides'].append(slide_data)

        return json_data

    except Exception as e:
        print(f"An error occurred: {str(e)}")


if __name__ == '__main__':
    # Parse the PowerPoint presentation and convert it into a JSON object
    json_data = parse_pptx_to_json('demo.pptx')

    # Save the JSON object to a file
    with open('presentation.json', 'w') as f:
        json.dump(json_data, f)
