import pptx
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ..pptx_handler.slide import Slide
from ..pptx_handler.slide_element import SlideElement
from ..utils.image_handler import ImageHandler
from ..utils.utils import extract_clickable_text
from ...config import EMU_CONVERSION_FACTOR, DIMENSION_PRECISION

def round_dimension(value):
    return round(value / EMU_CONVERSION_FACTOR, DIMENSION_PRECISION)

def extract_presentation_info(pptx_file, prs):
    presentation_name = os.path.basename(pptx_file).split('.')[0]
    presentation_width = round_dimension(prs.slide_width)
    presentation_height = round_dimension(prs.slide_height)

    return {
        'Module': presentation_name,
        'dimensions': {
            'width': presentation_width,
            'height': presentation_height
        }
    }

def get_text_properties(shape):
    paragraphs = shape.text_frame.paragraphs
    first_run = paragraphs[0].runs[0] if paragraphs and paragraphs[0].runs else None
    if not first_run:
        return {}
    
    return {
        'font_name': first_run.font.name,
        'font_size': first_run.font.size.pt if first_run.font.size else '',
        'font_bold': first_run.font.bold,
        'font_italic': first_run.font.italic,
        'font_underline': first_run.font.underline,
        'font_color': first_run.font.color.rgb if first_run.font.color and first_run.font.color.type == 'RGB' else '',
    }

def get_image_properties(shape, image_handler, slide_num):
    image_path = image_handler.save_image(shape.image.blob, slide_num, shape.name)
    return {
        'image_path': image_path,
        'image_properties': {
            'name': shape.name,
            'description': shape.description if hasattr(shape, 'description') else '',
            'extension': shape.image.ext,
            'content_type': shape.image.content_type,
            'file_size': shape.image.size,
            'dpi': shape.image.dpi,
        }
    }

def get_position_properties(shape):
    return {
        'left': round_dimension(shape.left),
        'top': round_dimension(shape.top),
        'width': round_dimension(shape.width),
        'height': round_dimension(shape.height)
    }

def process_slide(slide, image_handler, slide_num):
    slide_data = Slide(slide_num, title=slide.shapes.title.text if slide.shapes.title else None)

    if slide.has_notes_slide:
        slide_data.notes = slide.notes_slide.notes_text_frame.text
        slide_data.clickable_text = extract_clickable_text(slide_data.notes)

    for shape in slide.shapes:
        element_type = shape.shape_type.name if isinstance(shape.shape_type, MSO_SHAPE_TYPE) else shape.shape_type
        element = SlideElement(element_type)

        if hasattr(shape, 'text'):
            element.text = shape.text
            element.text_properties = get_text_properties(shape)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_properties = get_image_properties(shape, image_handler, slide_num)
            element.image_path = image_properties['image_path']
            element.image_properties = image_properties['image_properties']
        
        element.position = get_position_properties(shape)

        slide_data.add_element(element)

    return slide_data

def parse_pptx_to_json(pptx_file):
    prs = pptx.Presentation(pptx_file)
    data = extract_presentation_info(pptx_file, prs)
    image_handler = ImageHandler()

    data['slides'] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_obj = process_slide(slide, image_handler, slide_num)
        print(f"Processed slide {slide_num} with {len(slide_obj.elements)} elements.")
        data['slides'].append(slide_obj)

    return data
