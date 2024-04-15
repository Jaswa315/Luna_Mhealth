import pptx
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_handler.slide import Slide
from pptx_handler.slide_element import SlideElement
from utils.image_handler import ImageHandler
from utils.utils import extract_clickable_text

def extract_presentation_info(pptx_file, prs):
    presentation_name = os.path.basename(pptx_file).split('.')[0]  
    presentation_width = round(prs.slide_width / 360000, 4)
    presentation_height = round(prs.slide_height / 360000, 4)

    return {
        'Module': presentation_name,
        'dimensions': {
            'width': presentation_width,
            'height': presentation_height
        }
    }

def process_slide(slide, image_handler, slide_num):
    slide_data = Slide(slide_num, title=slide.shapes.title.text if slide.shapes.title else None)

    if slide.has_notes_slide:
        slide_data.notes = slide.notes_slide.notes_text_frame.text
        slide_data.clickable_text = extract_clickable_text(slide_data.notes)

    for shape in slide.shapes:
        element = SlideElement(shape.shape_type.name if isinstance(shape.shape_type, MSO_SHAPE_TYPE) else shape.shape_type)

        if hasattr(shape, 'text'):
            element.text = shape.text
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                element.text_properties = {
                    'font_name': shape.text_frame.paragraphs[0].runs[0].font.name,
                    'font_size': shape.text_frame.paragraphs[0].runs[0].font.size.pt if shape.text_frame.paragraphs[0].runs[0].font.size else '',
                    'font_bold': shape.text_frame.paragraphs[0].runs[0].font.bold,
                    'font_italic': shape.text_frame.paragraphs[0].runs[0].font.italic,
                    'font_underline': shape.text_frame.paragraphs[0].runs[0].font.underline,
                    'font_color': shape.text_frame.paragraphs[0].runs[0].font.color.rgb if shape.text_frame.paragraphs[0].runs[0].font.color.type == 'RGB' else '',
                }

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_path = image_handler.save_image(shape.image.blob, slide_num, shape.name)
            element.image_path = image_path
            element.image_properties = {
                        'name': shape.name,
                        'description': shape.description if hasattr(shape, 'description') else '',
                        'extension': shape.image.ext,
                        'content_type': shape.image.content_type,
                        'file_size': shape.image.size,
                        'dpi': shape.image.dpi,
                    }
        
        element.position = {
            'left': round(shape.left / 360000, 4),
            'top': round(shape.top / 360000, 4),
            'width': round(shape.width / 360000, 4),
            'height': round(shape.height / 360000, 4)
        }

        slide_data.add_element(element)

    return slide_data

def parse_pptx_to_json(pptx_file):
    prs = pptx.Presentation(pptx_file)
    data = extract_presentation_info(pptx_file, prs)
    image_handler = ImageHandler()

    data['slides'] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_obj = process_slide(slide, image_handler, slide_num)
        print(f"Processed slide {slide_num} with {len(slide_obj.elements)} elements.")
        data['slides'].append(slide_obj)

    return data 
